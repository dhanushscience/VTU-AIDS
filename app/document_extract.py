"""Extract plain text from uploaded reference documents."""

from __future__ import annotations

import io
from pathlib import Path

MAX_UPLOAD_BYTES = 25 * 1024 * 1024
MAX_CODE_BYTES = 512 * 1024
MAX_EXTRACTED_CHARS = 32_000

DOC_EXTENSIONS = frozenset({".pdf", ".pptx", ".docx", ".md", ".txt"})

IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"})

CODE_EXTENSIONS = frozenset(
    {
        ".py",
        ".js",
        ".ts",
        ".tsx",
        ".jsx",
        ".java",
        ".c",
        ".cc",
        ".cpp",
        ".cxx",
        ".h",
        ".hpp",
        ".cs",
        ".go",
        ".rs",
        ".rb",
        ".php",
        ".swift",
        ".kt",
        ".kts",
        ".scala",
        ".r",
        ".sql",
        ".html",
        ".htm",
        ".css",
        ".scss",
        ".sass",
        ".less",
        ".json",
        ".xml",
        ".yaml",
        ".yml",
        ".toml",
        ".ini",
        ".cfg",
        ".conf",
        ".sh",
        ".bash",
        ".zsh",
        ".ps1",
        ".bat",
        ".cmd",
        ".lua",
        ".vue",
        ".svelte",
        ".gradle",
        ".cmake",
        ".proto",
        ".ipynb",
    }
)

ALLOWED_EXTENSIONS = DOC_EXTENSIONS | IMAGE_EXTENSIONS | CODE_EXTENSIONS

EXT_LABELS: dict[str, str] = {
    ".pdf": "PDF",
    ".pptx": "PowerPoint",
    ".docx": "Word",
    ".md": "Markdown",
    ".txt": "Text",
    ".png": "Image",
    ".jpg": "Image",
    ".jpeg": "Image",
    ".gif": "Image",
    ".webp": "Image",
    ".bmp": "Image",
}

IMAGE_MIME = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
}

IMAGE_EXTRACT_PROMPT = """You are helping extract reference material for a VTU engineering internship diary (VTU AIDS).

From this image:
1. Transcribe all readable text (labels, slides, screenshots, handwriting if clear).
2. Summarize technical content: tools, hardware, code, diagrams, tasks, metrics.
3. List anything useful for writing daily internship log entries.

Return plain text only (no markdown code fences). Be factual and concise."""


def _clip(text: str) -> str:
    text = text.replace("\x00", " ").strip()
    if len(text) > MAX_EXTRACTED_CHARS:
        return text[:MAX_EXTRACTED_CHARS] + "\n\n[…truncated for length…]"
    return text


def _extract_pdf(data: bytes, *, api_key: str | None = None) -> str:
    reader = None
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        parts: list[str] = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                parts.append("")
        text = "\n".join(parts).strip()
        if text:
            return text
    except ModuleNotFoundError:
        # Runtime may not have pypdf (e.g. packaged env). Continue with Gemini fallback.
        pass
    except Exception:
        # Corrupt/unsupported PDF for local parser. Continue with Gemini fallback.
        pass

    # Fallback 1: direct PDF understanding via Gemini.
    if api_key:
        try:
            via_pdf = _extract_pdf_gemini(data, api_key)
            if via_pdf:
                return via_pdf
        except Exception:
            pass

    # Fallback 2: scanned/image-only PDFs via embedded image OCR.
    # Requires both API key and a parsed reader with embedded images.
    if not api_key or reader is None:
        return ""
    ocr_parts: list[str] = []
    for page in reader.pages[:8]:
        images = getattr(page, "images", None)
        if not images:
            continue
        for img in images[:2]:
            try:
                raw = getattr(img, "data", b"") or b""
                name = str(getattr(img, "name", "") or "").lower()
                if not raw:
                    continue
                if name.endswith(".png"):
                    mime = "image/png"
                elif name.endswith(".webp"):
                    mime = "image/webp"
                elif name.endswith(".gif"):
                    mime = "image/gif"
                elif name.endswith(".bmp"):
                    mime = "image/bmp"
                else:
                    mime = "image/jpeg"
                ocr = _extract_image_gemini(raw, mime, api_key).strip()
                if ocr:
                    ocr_parts.append(ocr)
            except Exception:
                continue
    return "\n\n".join(ocr_parts).strip()


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def _extract_plain(data: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_code(filename: str, data: bytes) -> str:
    if len(data) > MAX_CODE_BYTES:
        raise ValueError(
            f"Code file too large (max {MAX_CODE_BYTES // 1024} KB). Split or shorten the file."
        )
    text = _extract_plain(data).strip()
    if not text:
        raise ValueError("Code file is empty.")
    name = Path(filename).name
    ext = Path(filename).suffix.lower() or "text"
    return f"Source code file: {name}\n```{ext.lstrip('.')}\n{text}\n```"


def _extract_image_gemini(data: bytes, mime: str, api_key: str) -> str:
    from google import genai

    from app.config_store import normalize_api_key, validate_api_key_format
    from app.gemini_service import (
        DEFAULT_GEMINI_MODEL,
        generate_content_with_fallback,
        models_to_try,
        normalize_model_name,
    )
    from google.genai import types

    api_key = normalize_api_key(api_key)
    validate_api_key_format(api_key)

    client = genai.Client(api_key=api_key)
    parts = [
        types.Part.from_text(text=IMAGE_EXTRACT_PROMPT),
        types.Part.from_bytes(data=data, mime_type=mime),
    ]

    _, text = generate_content_with_fallback(
        client,
        models_to_try(normalize_model_name(DEFAULT_GEMINI_MODEL)),
        parts,
    )
    if text.strip():
        return text.strip()
    raise ValueError("Could not read image with Gemini (no text returned).")


def _extract_pdf_gemini(data: bytes, api_key: str) -> str:
    from google import genai
    from google.genai import types

    from app.config_store import normalize_api_key, validate_api_key_format
    from app.gemini_service import (
        DEFAULT_GEMINI_MODEL,
        generate_content_with_fallback,
        models_to_try,
        normalize_model_name,
    )

    api_key = normalize_api_key(api_key)
    validate_api_key_format(api_key)
    client = genai.Client(api_key=api_key)
    parts = [
        types.Part.from_text(
            text=(
                "Extract readable text and key technical points from this PDF for internship diary writing. "
                "Return plain text only."
            )
        ),
        types.Part.from_bytes(data=data, mime_type="application/pdf"),
    ]

    _, text = generate_content_with_fallback(
        client,
        models_to_try(normalize_model_name(DEFAULT_GEMINI_MODEL)),
        parts,
    )
    return text.strip()


def _kind_label(ext: str) -> str:
    if ext in IMAGE_EXTENSIONS:
        return EXT_LABELS.get(ext, "Image")
    if ext in CODE_EXTENSIONS:
        return "Code"
    return EXT_LABELS.get(ext, ext.lstrip(".").upper() or "File")


def extract_text_from_upload(
    filename: str,
    data: bytes,
    *,
    api_key: str | None = None,
) -> dict[str, str | int]:
    """Return {filename, kind, text, char_count}. Raises ValueError on bad input."""
    if len(data) > MAX_UPLOAD_BYTES:
        raise ValueError(f"File too large (max {MAX_UPLOAD_BYTES // (1024 * 1024)} MB).")

    ext = Path(filename or "").suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(
            "Unsupported file type. Use PDF, PPTX, DOCX, images (PNG/JPG/…), "
            "code files (.py, .js, …), Markdown, or TXT."
        )

    if ext in IMAGE_EXTENSIONS:
        if not api_key:
            raise ValueError(
                "Images need a Gemini API key. Add your key in Settings, then upload again."
            )
        mime = IMAGE_MIME.get(ext, "image/jpeg")
        text = _extract_image_gemini(data, mime, api_key)
    elif ext == ".pdf":
        text = _extract_pdf(data, api_key=api_key)
    elif ext == ".docx":
        text = _extract_docx(data)
    elif ext == ".doc":
        try:
            text = _extract_docx(data)
        except Exception as e:
            raise ValueError(
                "Old .doc format is not supported. Save the file as .docx and upload again."
            ) from e
    elif ext == ".pptx":
        text = _extract_pptx(data)
    elif ext in CODE_EXTENSIONS:
        text = _extract_code(filename, data)
    else:
        text = _extract_plain(data)

    text = _clip(text)
    if not text:
        raise ValueError("No readable text found in this file.")

    return {
        "filename": Path(filename).name,
        "kind": _kind_label(ext),
        "text": text,
        "char_count": len(text),
    }
